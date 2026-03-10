from __future__ import annotations

import re
import shutil
from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader, PdfWriter
from pypdf.generic import BooleanObject, NameObject, TextStringObject
from docx import Document
from pptx import Presentation


def detect_type(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".docx":
        return "docx"
    if ext == ".pptx":
        return "pptx"
    if ext == ".pdf":
        return "pdf"
    if ext in {".txt", ".md"}:
        return "text"
    raise ValueError(f"Unsupported file type: {path}")


BLANK_PATTERNS = [
    r"_{3,}",
    r"-{3,}",
    r"[ \t]{6,}",
    r"[\u2000-\u200A]{3,}",
    r"[\u00A0]{3,}",
]
BLANK_REGEX = re.compile("|".join(f"(?:{p})" for p in BLANK_PATTERNS))


@dataclass
class DocxBlankSlot:
    blank_id: int
    para_index: int
    start: int
    end: int
    context: str
    label_hint: str
    nearby_labels: list[str]


@dataclass
class PptxBlankSlot:
    blank_id: int
    slide_index: int
    shape_index: int
    start: int
    end: int
    context: str
    label_hint: str


def extract_text(path: Path) -> str:
    t = detect_type(path)
    if t == "docx":
        return extract_docx_text(path)
    if t == "pptx":
        return extract_pptx_text(path)
    if t == "pdf":
        return extract_pdf_text(path)
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_docx_text(path: Path) -> str:
    doc = Document(path)
    lines: list[str] = []
    for p in _iter_docx_paragraphs(doc):
        if p.text:
            lines.append(p.text)
    return "\n".join(lines)


def extract_pptx_text(path: Path) -> str:
    prs = Presentation(path)
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def _iter_docx_paragraphs(doc: Document):
    for p in doc.paragraphs:
        yield p

    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    yield p

    for section in doc.sections:
        for p in section.header.paragraphs:
            yield p
        for t in section.header.tables:
            for row in t.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        yield p
        for p in section.footer.paragraphs:
            yield p
        for t in section.footer.tables:
            for row in t.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        yield p


def _paragraph_full_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def _set_paragraph_text_single_run(paragraph, text: str) -> None:
    if not paragraph.runs:
        paragraph.add_run(text)
        return
    for run in paragraph.runs:
        run.text = ""
    paragraph.runs[0].text = text


def _derive_label_hint(full_text: str, blank_start: int) -> str:
    left = full_text[:blank_start].replace("\n", " ").strip()
    if not left:
        return ""
    tail = left[-120:]
    # Prefer the most recent labeled segment before the blank.
    m = re.search(r"([A-Za-z][A-Za-z0-9/&() ,'\-]{2,60})\s*[:\-–—]?\s*$", tail)
    if not m:
        return ""
    return m.group(1).strip(" -:;,.")


def _derive_label_hint_from_previous_paragraphs(paragraphs, para_index: int, lookback: int = 6) -> str:
    start = max(0, para_index - lookback)
    for i in range(para_index - 1, start - 1, -1):
        t = _paragraph_full_text(paragraphs[i]).replace("\n", " ").strip()
        if not t:
            continue
        if BLANK_REGEX.fullmatch(t):
            continue
        if t.upper() == "FOR THE WEEK":
            continue
        if t.startswith("(") and t.endswith(")"):
            continue
        if len(t) > 100:
            continue
        # Reuse the same cleanup logic as local hint derivation.
        m = re.search(r"([A-Za-z][A-Za-z0-9/&() ,'\-]{2,80})\s*[:\-–—]?\s*$", t)
        if not m:
            continue
        return m.group(1).strip(" -:;,.")
    return ""


def _collect_previous_labels(paragraphs, para_index: int, lookback: int = 12) -> list[str]:
    labels: list[str] = []
    start = max(0, para_index - lookback)
    for i in range(para_index - 1, start - 1, -1):
        t = _paragraph_full_text(paragraphs[i]).replace("\n", " ").strip()
        if not t:
            continue
        if BLANK_REGEX.fullmatch(t):
            continue
        if t.startswith("(") and t.endswith(")"):
            continue
        if len(t) > 100:
            continue
        labels.append(t.strip(" -:;,."))
        if len(labels) >= 6:
            break
    return labels


def collect_docx_blanks(path: Path, max_blanks: int = 120, context_window: int = 60) -> list[DocxBlankSlot]:
    doc = Document(path)
    blanks: list[DocxBlankSlot] = []
    paragraphs = list(_iter_docx_paragraphs(doc))
    for para_index, paragraph in enumerate(paragraphs):
        full = _paragraph_full_text(paragraph)
        for match in BLANK_REGEX.finditer(full):
            s, e = match.start(), match.end()
            left = max(0, s - context_window)
            right = min(len(full), e + context_window)
            blanks.append(
                DocxBlankSlot(
                    blank_id=len(blanks),
                    para_index=para_index,
                    start=s,
                    end=e,
                    context=full[left:right].replace("\n", " "),
                    label_hint=_derive_label_hint(full, s)
                    or _derive_label_hint_from_previous_paragraphs(paragraphs, para_index),
                    nearby_labels=_collect_previous_labels(paragraphs, para_index),
                )
            )
            if len(blanks) >= max_blanks:
                return blanks
    return blanks


def apply_docx_blank_placeholders(
    source: Path,
    destination: Path,
    blank_to_placeholder: dict[int, str],
    blanks: list[DocxBlankSlot],
) -> None:
    doc = Document(source)
    paragraphs = list(_iter_docx_paragraphs(doc))

    by_para: dict[int, list[tuple[int, int, str]]] = {}
    for slot in blanks:
        placeholder = blank_to_placeholder.get(slot.blank_id)
        if not placeholder:
            continue
        by_para.setdefault(slot.para_index, []).append((slot.start, slot.end, placeholder))

    for para_index, reps in by_para.items():
        paragraph = paragraphs[para_index]
        full = _paragraph_full_text(paragraph)
        new_text = full
        for s, e, placeholder in sorted(reps, key=lambda x: x[0], reverse=True):
            new_text = new_text[:s] + placeholder + new_text[e:]
        if new_text != full:
            _set_paragraph_text_single_run(paragraph, new_text)

    doc.save(destination)


def collect_pptx_blanks(path: Path, max_blanks: int = 120, context_window: int = 60) -> list[PptxBlankSlot]:
    prs = Presentation(path)
    blanks: list[PptxBlankSlot] = []
    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if not hasattr(shape, "text"):
                continue
            full = shape.text or ""
            for match in BLANK_REGEX.finditer(full):
                s, e = match.start(), match.end()
                left = max(0, s - context_window)
                right = min(len(full), e + context_window)
                blanks.append(
                    PptxBlankSlot(
                        blank_id=len(blanks),
                        slide_index=slide_index,
                        shape_index=shape_index,
                        start=s,
                        end=e,
                        context=full[left:right].replace("\n", " "),
                        label_hint=_derive_label_hint(full, s),
                    )
                )
                if len(blanks) >= max_blanks:
                    return blanks
    return blanks


def apply_pptx_blank_placeholders(
    source: Path,
    destination: Path,
    blank_to_placeholder: dict[int, str],
    blanks: list[PptxBlankSlot],
) -> None:
    prs = Presentation(source)

    by_shape: dict[tuple[int, int], list[tuple[int, int, str]]] = {}
    for slot in blanks:
        placeholder = blank_to_placeholder.get(slot.blank_id)
        if not placeholder:
            continue
        key = (slot.slide_index, slot.shape_index)
        by_shape.setdefault(key, []).append((slot.start, slot.end, placeholder))

    for (slide_index, shape_index), reps in by_shape.items():
        shape = prs.slides[slide_index].shapes[shape_index]
        full = shape.text or ""
        new_text = full
        for s, e, placeholder in sorted(reps, key=lambda x: x[0], reverse=True):
            new_text = new_text[:s] + placeholder + new_text[e:]
        if new_text != full:
            shape.text = new_text

    prs.save(destination)


def extract_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    lines: list[str] = []
    for page in reader.pages:
        lines.append(page.extract_text() or "")
    return "\n".join(lines)


def list_pdf_form_fields(path: Path) -> list[str]:
    reader = PdfReader(str(path))
    fields = reader.get_fields()
    if not fields:
        return []
    names: list[str] = []
    for name, field in fields.items():
        if not field:
            continue
        if field.get("/FT") is None:
            continue
        names.append(str(name))
    return sorted(set(names))


def _extract_on_state(field: dict) -> str:
    states = field.get("/_States_")
    if states:
        for state in states:
            s = str(state)
            if s != "/Off":
                return s
    ap = field.get("/AP")
    if not ap:
        return "/Yes"
    normal = ap.get("/N") if hasattr(ap, "get") else None
    if not normal:
        return "/Yes"
    keys = [str(k) for k in normal.keys()]
    for k in keys:
        if k != "/Off":
            return k
    return "/Yes"


def _resolve_inherited(widget: dict, key: str):
    current = widget
    for _ in range(24):
        if key in current:
            return current[key]
        parent = current.get("/Parent")
        if not parent:
            break
        current = parent.get_object()
    return None


def _qualified_widget_name(widget: dict) -> str:
    parts: list[str] = []
    current = widget
    for _ in range(24):
        name = current.get("/T")
        if name:
            parts.append(str(name))
        parent = current.get("/Parent")
        if not parent:
            break
        current = parent.get_object()
    return ".".join(reversed(parts))


def _set_widget_value(widget: dict, value: object) -> None:
    widget[NameObject("/V")] = value
    parent = widget.get("/Parent")
    if parent:
        p = parent.get_object()
        p[NameObject("/V")] = value


def _normalize_pdf_value(field_name: str, field: dict, raw: str) -> object:
    ft = str(field.get("/FT") or "")
    text = str(raw if raw is not None else "").strip()
    if ft == "/Btn":
        lowered = text.lower()
        on_state = _extract_on_state(field)
        if text.startswith("/") and text in {"/Off", on_state}:
            return NameObject(text)
        if text and ("/" + text) in {"/Off", on_state}:
            return NameObject("/" + text)
        if lowered in {"1", "true", "yes", "y", "checked", "check", "on", "x"}:
            return NameObject(on_state)
        if lowered in {"0", "false", "no", "n", "unchecked", "off", ""}:
            return NameObject("/Off")
        return NameObject(on_state)
    if ft in {"/Tx", "/Ch"}:
        return text
    return text


def replace_in_docx(source: Path, destination: Path, replacements: dict[str, str]) -> None:
    doc = Document(source)
    active_replacements = {old: new for old, new in replacements.items() if old}

    for paragraph in _iter_docx_paragraphs(doc):
        for run in paragraph.runs:
            original = run.text
            text = original
            for old, new in active_replacements.items():
                if old in text:
                    text = text.replace(old, new)
            # Only rewrite runs that actually changed. Rewriting unchanged runs
            # can remove non-text run internals (for example legacy form fields).
            if text != original:
                run.text = text

        # Fallback for placeholders split across multiple runs.
        full_after_run_pass = _paragraph_full_text(paragraph)
        merged = full_after_run_pass
        for old, new in active_replacements.items():
            if old in merged:
                merged = merged.replace(old, new)
        if merged != full_after_run_pass:
            _set_paragraph_text_single_run(paragraph, merged)

    doc.save(destination)


def replace_in_pptx(source: Path, destination: Path, replacements: dict[str, str]) -> None:
    prs = Presentation(source)
    active_replacements = {old: new for old, new in replacements.items() if old}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            original = shape.text
            text = original
            for old, new in active_replacements.items():
                if old in text:
                    text = text.replace(old, new)
            if text != original:
                shape.text = text
    prs.save(destination)


def replace_in_text(source: Path, destination: Path, replacements: dict[str, str]) -> None:
    text = source.read_text(encoding="utf-8", errors="ignore")
    for old, new in replacements.items():
        if old:
            text = text.replace(old, new)
    destination.write_text(text, encoding="utf-8")


def fill_pdf_form(source: Path, destination: Path, values: dict[str, str]) -> None:
    reader = PdfReader(str(source))
    writer = PdfWriter()

    # Clone full reader structure first so AcroForm metadata is preserved.
    writer.clone_document_from_reader(reader)

    # Some PDFs still miss /AcroForm in writer root after clone; copy it explicitly.
    reader_root = reader.trailer.get("/Root", {})
    writer_root = writer._root_object
    if "/AcroForm" not in writer_root and "/AcroForm" in reader_root:
        writer_root[NameObject("/AcroForm")] = reader_root["/AcroForm"]
    acro = writer_root.get("/AcroForm")
    if acro:
        acro_obj = acro.get_object()
        if "/XFA" in acro_obj:
            # XFA often overrides AcroForm values and causes viewers to show fields as blank.
            del acro_obj[NameObject("/XFA")]
        acro_obj[NameObject("/NeedAppearances")] = BooleanObject(True)

    fields = reader.get_fields() or {}
    normalized_values: dict[str, object] = {}
    for field_name, field in fields.items():
        if not field or field.get("/FT") is None:
            continue
        if field_name not in values:
            continue
        normalized_values[field_name] = _normalize_pdf_value(field_name, field, values[field_name])

    # Fallback for callers that already provide exact field names without field metadata.
    if not normalized_values:
        normalized_values = {str(k): str(v) for k, v in values.items()}

    # Fill widgets by fully-qualified field path to avoid writing to similarly named fields.
    for page in writer.pages:
        annots = page.get("/Annots") or []
        for annot_ref in annots:
            annot = annot_ref.get_object()
            if str(annot.get("/Subtype", "")) != "/Widget":
                continue
            full_name = _qualified_widget_name(annot)
            short_name = str(annot.get("/T", ""))
            if full_name in normalized_values:
                normalized = normalized_values[full_name]
            elif short_name in normalized_values:
                normalized = normalized_values[short_name]
            else:
                continue

            ft = str(_resolve_inherited(annot, "/FT") or "")
            if ft == "/Btn":
                if isinstance(normalized, NameObject):
                    chosen = normalized
                else:
                    chosen = NameObject(str(normalized))
                on_state = _extract_on_state(annot)
                if str(chosen) != "/Off" and str(chosen) != on_state:
                    chosen = NameObject(on_state)
                _set_widget_value(annot, chosen)
                annot[NameObject("/AS")] = chosen
            else:
                txt = normalized if isinstance(normalized, TextStringObject) else TextStringObject(str(normalized))
                _set_widget_value(annot, txt)

    with destination.open("wb") as f:
        writer.write(f)


def copy_file(source: Path, destination: Path) -> None:
    shutil.copy2(source, destination)
